"""PowerPoint converter with slide and heading rules."""

from __future__ import annotations

import re

from pptx import Presentation
from pptx.util import Inches, Pt

from shared.converter_interface import DocumentConverter
from shared.markdown_parser import SLIDE_RE, parse_table_block
from shared.models import ConversionJob


def _add_bullets(text_frame, lines: list[str]) -> None:
    text_frame.clear()
    first = True
    for raw in lines:
        line = raw.strip()
        if not line or line.startswith("|") or line.startswith("```"):
            continue
        if line.startswith("- "):
            line = line[2:]
        paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.font.size = Pt(18)
        first = False


def _add_content_slide(prs: Presentation, title: str, chunk: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    bullet_lines = [ln for ln in chunk.splitlines() if ln.strip().startswith("- ")]
    if not bullet_lines:
        bullet_lines = [
            ln for ln in chunk.splitlines() if ln.strip() and not ln.startswith("#") and not ln.startswith("|")
        ]
    _add_bullets(slide.placeholders[1].text_frame, bullet_lines[:12])
    table_rows = parse_table_block([line for line in chunk.splitlines() if line.strip().startswith("|")])
    if table_rows and len(table_rows) >= 2:
        rows, cols = len(table_rows), len(table_rows[0])
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(4.0), Inches(9), Inches(2.5))
        table = table_shape.table
        for r, row in enumerate(table_rows):
            for c, val in enumerate(row):
                table.cell(r, c).text = val


class PPTConverter(DocumentConverter):
    format_ext = ".pptx"
    format_name = "pptx"

    def convert(self, job: ConversionJob) -> None:
        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = job.meta.get("title", "Presentation")
        title_slide.placeholders[1].text = (
            f"{job.meta.get('project', '')} | {job.meta.get('last_updated', '')}"
        )

        slide_matches = list(SLIDE_RE.finditer(job.body))
        if slide_matches:
            for idx, match in enumerate(slide_matches):
                slide_title = match.group(2).strip()
                start = match.end()
                end = slide_matches[idx + 1].start() if idx + 1 < len(slide_matches) else len(job.body)
                _add_content_slide(prs, slide_title, job.body[start:end])
        else:
            h1_parts = re.split(r"^#\s+(.+)$", job.body, flags=re.MULTILINE)
            if len(h1_parts) > 1:
                i = 1
                while i < len(h1_parts):
                    section_title = h1_parts[i].strip()
                    section_body = h1_parts[i + 1] if i + 1 < len(h1_parts) else ""
                    prs.slides.add_slide(prs.slide_layouts[2]).shapes.title.text = section_title
                    for h2_match in re.finditer(r"^##\s+(.+)$", section_body, re.MULTILINE):
                        h2_title = h2_match.group(1).strip()
                        h2_start = h2_match.end()
                        h2_next = re.search(r"^##\s+", section_body[h2_start:], re.MULTILINE)
                        h2_end = h2_start + h2_next.start() if h2_next else len(section_body)
                        _add_content_slide(prs, h2_title, section_body[h2_start:h2_end])
                    i += 2
            else:
                for h2_match in re.finditer(r"^##\s+(.+)$", job.body, re.MULTILINE):
                    h2_title = h2_match.group(1).strip()
                    if h2_title.lower().startswith("slide"):
                        continue
                    h2_start = h2_match.end()
                    h2_next = re.search(r"^##\s+", job.body[h2_start:], re.MULTILINE)
                    h2_end = h2_start + h2_next.start() if h2_next else len(job.body)
                    _add_content_slide(prs, h2_title, job.body[h2_start:h2_end])

        body_lower = job.body.lower()
        if "thank you" not in body_lower:
            thank = prs.slides.add_slide(prs.slide_layouts[0])
            thank.shapes.title.text = "Thank You"
            thank.placeholders[1].text = job.meta.get("project", "")

        prs.save(job.output_path)
